from pptx import Presentation

def extract_pptx(path):
    pres = Presentation(path)
    texts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

